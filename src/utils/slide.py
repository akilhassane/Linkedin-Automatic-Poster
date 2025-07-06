from pathlib import Path
from typing import List
from pptx import Presentation  # type: ignore
from pptx.util import Inches, Pt  # type: ignore
import logging


def build_slide(topic: str, bullet_points: List[str], image_path: str | Path | None = None) -> str:
    """Creates a single-slide PPTX summarising bullet points and optional image."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = f"Daily Digest: {topic.title()}"

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.level = 0

    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(5), height=Inches(2))

    out_path = f"digest_{topic.replace(' ', '_')}.pptx"
    prs.save(out_path)
    logging.info("Slide saved at %s", out_path)
    return out_path