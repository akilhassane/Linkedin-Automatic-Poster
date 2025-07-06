"""generator.py

Creates the visual and textual artefacts that will be attached to the
LinkedIn post.

Outputs
-------
1. ``article_body``  – plain-text LinkedIn post body (<= 1 300 chars).
2. ``slide_path``    – path to a temporary PowerPoint file (*.pptx*) containing
                       a title slide + summary slides.
3. ``graph_path``    – path to a temporary PNG file with a simple bar/line
                       chart derived from numeric snippets in the articles.

Public helper: ``build_assets(raw_articles, summaries)``

Design notes
------------
• Slide deck uses *python-pptx* (blank theme) so that LinkedIn treats it as
a native document. A title slide is followed by up to ``config.SLIDES_MAX``
content slides (title + bullet points).
• Numeric extraction is intentionally simple (regex) – enough to produce a
quick visual. If no numbers are found, we chart article lengths (word counts)
as a fallback.
"""
from __future__ import annotations

import datetime as _dt
import os
import pathlib
import re
import textwrap
from typing import List, Dict, Tuple, Any

import matplotlib.pyplot as _plt  # type: ignore
from pptx import Presentation  # type: ignore
from pptx.util import Inches, Pt  # type: ignore

import config

_ASSETS_DIR = pathlib.Path("assets")
_ASSETS_DIR.mkdir(exist_ok=True)

# ----------------------------
# Post body builder
# ----------------------------


def _make_post_body(raw_articles: List[Dict[str, Any]], summaries: List[str]) -> str:
    """Compose a LinkedIn-friendly body text.

    We concatenate each article title + summary into a single paragraph. We
    aim to stay under 1 300 characters – LinkedIn's max for regular posts –
    but also provide a "Read more" link for each source.
    """
    parts: List[str] = []
    for art, summ in zip(raw_articles, summaries):
        title = art.get("title", "Untitled")
        url = art.get("url", "")
        blob = f"• {title}: {summ}".strip()
        if url:
            blob += f" (source: {url})"
        parts.append(blob)

    body = "\n\n".join(parts)

    # Truncate if necessary.
    if len(body) > 1280:  # leave leeway for hashtags/footer.
        body = body[:1270].rsplit(" ", 1)[0] + " …"

    footer = f"\n\n#AI #automation\nPosted on {_dt.datetime.utcnow().strftime('%Y-%m-%d')} via bot"
    return body + footer


# ----------------------------
# Slide deck builder
# ----------------------------

def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle  # type: ignore[index]


def _add_summary_slide(prs: Presentation, heading: str, summary: str) -> None:
    layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = heading

    tx_box = slide.shapes.placeholders[1].text_frame  # type: ignore[index]
    tx_box.clear()

    # Wrap summary into bullet points of ~20 words.
    bullets = textwrap.wrap(summary, width=100)
    for idx, bullet in enumerate(bullets):
        p = tx_box.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(14)
        if idx == 0:
            # First paragraph keeps bullet formatting already set.
            first = tx_box.paragraphs[0]
            first.text = bullet
            first.level = 0
            first.font.size = Pt(14)
            tx_box._p.remove(tx_box._p_lst[1])  # remove redundant paragraph
            break


def _build_pptx(raw_articles: List[Dict[str, Any]], summaries: List[str]) -> str:
    prs = Presentation()
    today = _dt.datetime.utcnow().strftime("%Y-%m-%d")
    _add_title_slide(prs, f"{config.TOPIC.title()} – Daily Digest", today)

    for art, summ in list(zip(raw_articles, summaries))[: config.SLIDES_MAX]:
        _add_summary_slide(prs, art.get("title", "Untitled"), summ)

    fname = _ASSETS_DIR / f"slides_{today}.pptx"
    prs.save(fname)
    return str(fname)


# ----------------------------
# Chart builder
# ----------------------------

def _extract_numbers(text: str) -> List[float]:
    pattern = re.compile(r"[-+]?\d+(?:\.\d+)?")
    return [float(m.group()) for m in pattern.finditer(text)]


def _build_chart(raw_articles: List[Dict[str, Any]]) -> str:
    # Derive one numeric value per article (first found) else fallback to word count.
    xs, ys = [], []
    for idx, art in enumerate(raw_articles, 1):
        nums = _extract_numbers(art["text"])
        val = nums[0] if nums else len(art["text"].split())
        xs.append(f"Art {idx}")
        ys.append(val)

    _plt.figure(figsize=(6, 4))
    if any(_extract_numbers(art["text"]) for art in raw_articles):
        _plt.bar(xs, ys, color="skyblue")
        _plt.ylabel("Value (first numeric in article)")
    else:
        _plt.plot(xs, ys, marker="o")
        _plt.ylabel("Article length (words)")
    _plt.title(f"{config.TOPIC.title()} – Daily Metrics")
    _plt.tight_layout()

    fname = _ASSETS_DIR / f"chart_{_dt.datetime.utcnow().strftime('%Y-%m-%d')}.png"
    _plt.savefig(fname, dpi=150)
    _plt.close()
    return str(fname)


# ----------------------------
# Facade
# ----------------------------

def build_assets(
    raw_articles: List[Dict[str, Any]],
    summaries: List[str],
) -> Tuple[str, str, str]:
    """Return (post_body, slide_path, graph_path)."""
    body = _make_post_body(raw_articles, summaries)
    slide_path = _build_pptx(raw_articles, summaries)
    graph_path = _build_chart(raw_articles)
    return body, slide_path, graph_path


if __name__ == "__main__":
    # Simple manual test with dummy articles.
    dummy = [
        {"title": "Article 1", "text": "The market grew by 15% last year."},
        {"title": "Article 2", "text": "A total of 42 qubits were achieved in 2024."},
    ]
    dummy_summaries = [
        "Market grew by fifteen percent.",
        "Forty-two qubits achieved in 2024.",
    ]
    print(build_assets(dummy, dummy_summaries)[0])